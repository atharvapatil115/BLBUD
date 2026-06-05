import os
import pdfplumber
from docx import Document
from pptx import Presentation


class DocumentExtractor:
    """
    Extract text from PDF, DOCX and PPTX files.
    """

    def __init__(self):
        pass

    def extract_pdf(self, file_path):
        """
        Extract text from PDF file.
        """
        text = []

        try:
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, start=1):
                    page_text = page.extract_text()

                    if page_text:
                        text.append(
                            f"\n--- PAGE {page_num} ---\n{page_text}"
                        )

            return "\n".join(text)

        except Exception as e:
            print(f"PDF Extraction Error: {e}")
            return ""

    def extract_docx(self, file_path):
        """
        Extract text from DOCX file.
        """
        try:
            doc = Document(file_path)

            text = [
                para.text
                for para in doc.paragraphs
                if para.text.strip()
            ]

            return "\n".join(text)

        except Exception as e:
            print(f" DOCX Extraction Error: {e}")
            return ""

    def extract_pptx(self, file_path):
        """
        Extract text from PPTX file.
        """
        try:
            presentation = Presentation(file_path)

            text = []

            for slide_num, slide in enumerate(
                presentation.slides,
                start=1
            ):
                text.append(
                    f"\n--- SLIDE {slide_num} ---\n"
                )

                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            text.append(shape.text)

            return "\n".join(text)

        except Exception as e:
            print(f"PPTX Extraction Error: {e}")
            return ""

    def process_document(self, file_path):
        """
        Automatically detect file type
        and extract text.
        """

        extension = os.path.splitext(file_path)[1].lower()

        if extension == ".pdf":
            return self.extract_pdf(file_path)

        elif extension == ".docx":
            return self.extract_docx(file_path)

        elif extension == ".pptx":
            return self.extract_pptx(file_path)

        else:
            print(f"⚠ Unsupported file type: {extension}")
            return ""

    def save_text(self, text, output_path):
        """
        Save extracted text.
        """
        try:
            with open(
                output_path,
                "w",
                encoding="utf-8"
            ) as file:
                file.write(text)

            print(f"Saved: {output_path}")

        except Exception as e:
            print(f" Save Error: {e}")

    def process_folder(
        self,
        input_folder,
        output_folder
    ):
        """
        Process all documents inside folder.
        """

        os.makedirs(
            output_folder,
            exist_ok=True
        )

        supported_extensions = (
            ".pdf",
            ".docx",
            ".pptx"
        )

        for file_name in os.listdir(input_folder):

            if not file_name.lower().endswith(
                supported_extensions
            ):
                continue

            file_path = os.path.join(
                input_folder,
                file_name
            )

            print(f" Processing: {file_name}")

            extracted_text = self.process_document(
                file_path
            )

            output_file = (
                os.path.splitext(file_name)[0]
                + ".txt"
            )

            output_path = os.path.join(
                output_folder,
                output_file
            )

            self.save_text(
                extracted_text,
                output_path
            )

        print("\nDocument Extraction Completed")